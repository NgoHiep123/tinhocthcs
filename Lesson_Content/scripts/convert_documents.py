#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script chuyển đổi tài liệu (PDF, DOCX, PPTX) sang các định dạng web-friendly
"""

import os
import sys
from pathlib import Path
import subprocess
import json

class DocumentConverter:
    """Chuyển đổi tài liệu sang định dạng web"""
    
    def __init__(self, output_dir='converted'):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def convert_pdf_to_images(self, pdf_path, output_folder=None):
        """
        Chuyển PDF thành ảnh (mỗi trang một ảnh)
        Yêu cầu: pdf2image (pip install pdf2image) và poppler
        """
        try:
            from pdf2image import convert_from_path
            
            if output_folder is None:
                output_folder = self.output_dir / Path(pdf_path).stem
            
            output_folder = Path(output_folder)
            output_folder.mkdir(parents=True, exist_ok=True)
            
            print(f"🔄 Đang chuyển đổi PDF: {pdf_path}")
            
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=150)
            
            image_paths = []
            for i, image in enumerate(images, 1):
                image_path = output_folder / f"page_{i:03d}.jpg"
                image.save(image_path, 'JPEG', quality=85)
                image_paths.append(str(image_path))
                print(f"  ✓ Trang {i}/{len(images)}")
            
            print(f"✅ Đã chuyển đổi {len(images)} trang")
            return image_paths
            
        except ImportError:
            print("❌ Chưa cài đặt pdf2image. Cài đặt: pip install pdf2image")
            print("   Cần cài đặt poppler: https://github.com/oschwartz10612/poppler-windows/releases/")
            return []
        except Exception as e:
            print(f"❌ Lỗi chuyển đổi PDF: {e}")
            return []
    
    def convert_docx_to_html(self, docx_path, output_html=None):
        """
        Chuyển DOCX sang HTML
        Yêu cầu: python-docx (pip install python-docx)
        """
        try:
            from docx import Document
            from docx.shared import Inches
            
            if output_html is None:
                output_html = self.output_dir / f"{Path(docx_path).stem}.html"
            
            print(f"🔄 Đang chuyển đổi DOCX: {docx_path}")
            
            doc = Document(docx_path)
            
            # Tạo HTML
            html_parts = []
            html_parts.append('<!DOCTYPE html>')
            html_parts.append('<html lang="vi">')
            html_parts.append('<head>')
            html_parts.append('<meta charset="UTF-8">')
            html_parts.append('<meta name="viewport" content="width=device-width, initial-scale=1.0">')
            html_parts.append('<title>Document</title>')
            html_parts.append('<style>')
            html_parts.append('''
                body { font-family: 'Times New Roman', serif; line-height: 1.6; max-width: 800px; margin: 0 auto; padding: 20px; }
                h1 { font-size: 24px; font-weight: bold; margin: 20px 0 10px; }
                h2 { font-size: 20px; font-weight: bold; margin: 18px 0 8px; }
                h3 { font-size: 18px; font-weight: bold; margin: 16px 0 8px; }
                p { margin: 10px 0; text-align: justify; }
                ul, ol { margin: 10px 0; padding-left: 40px; }
                li { margin: 5px 0; }
                strong { font-weight: bold; }
                em { font-style: italic; }
                img { max-width: 100%; height: auto; margin: 10px 0; }
            ''')
            html_parts.append('</style>')
            html_parts.append('</head>')
            html_parts.append('<body>')
            
            # Convert paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name.lower()
                    
                    if 'heading 1' in style:
                        html_parts.append(f'<h1>{self._escape_html(para.text)}</h1>')
                    elif 'heading 2' in style:
                        html_parts.append(f'<h2>{self._escape_html(para.text)}</h2>')
                    elif 'heading 3' in style:
                        html_parts.append(f'<h3>{self._escape_html(para.text)}</h3>')
                    else:
                        html_parts.append(f'<p>{self._escape_html(para.text)}</p>')
            
            html_parts.append('</body>')
            html_parts.append('</html>')
            
            html_content = '\n'.join(html_parts)
            
            # Save HTML
            with open(output_html, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            print(f"✅ Đã chuyển đổi sang: {output_html}")
            return str(output_html)
            
        except ImportError:
            print("❌ Chưa cài đặt python-docx. Cài đặt: pip install python-docx")
            return None
        except Exception as e:
            print(f"❌ Lỗi chuyển đổi DOCX: {e}")
            return None
    
    def convert_pptx_to_images(self, pptx_path, output_folder=None):
        """
        Chuyển PPTX thành ảnh (mỗi slide một ảnh)
        Yêu cầu: python-pptx (pip install python-pptx) và LibreOffice/PowerPoint
        """
        try:
            from pptx import Presentation
            
            if output_folder is None:
                output_folder = self.output_dir / Path(pptx_path).stem
            
            output_folder = Path(output_folder)
            output_folder.mkdir(parents=True, exist_ok=True)
            
            print(f"🔄 Đang phân tích PPTX: {pptx_path}")
            
            prs = Presentation(pptx_path)
            
            print(f"📊 Tìm thấy {len(prs.slides)} slide")
            
            # Extract notes
            notes = []
            for i, slide in enumerate(prs.slides, 1):
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    notes.append(notes_text)
                else:
                    notes.append("")
            
            # Save notes
            notes_file = output_folder / 'notes.json'
            with open(notes_file, 'w', encoding='utf-8') as f:
                json.dump(notes, f, ensure_ascii=False, indent=2)
            
            print(f"✅ Đã lưu ghi chú: {notes_file}")
            
            # Convert to images using LibreOffice
            print("⚠️ Để chuyển PPTX sang ảnh, cần:")
            print("  1. LibreOffice (https://www.libreoffice.org/download/)")
            print("  2. Hoặc sử dụng PowerPoint để Export as Images")
            print(f"\n📝 Hướng dẫn thủ công:")
            print(f"  1. Mở file: {pptx_path}")
            print(f"  2. File > Export > Change File Type > JPEG")
            print(f"  3. Lưu vào: {output_folder}")
            
            return {
                'slides': [],  # Sẽ được điền sau khi convert thủ công
                'notes': notes,
                'total_slides': len(prs.slides),
                'notes_file': str(notes_file)
            }
            
        except ImportError:
            print("❌ Chưa cài đặt python-pptx. Cài đặt: pip install python-pptx")
            return None
        except Exception as e:
            print(f"❌ Lỗi xử lý PPTX: {e}")
            return None
    
    def _escape_html(self, text):
        """Escape HTML characters"""
        return (text
                .replace('&', '&amp;')
                .replace('<', '&lt;')
                .replace('>', '&gt;')
                .replace('"', '&quot;')
                .replace("'", '&#39;'))
    
    def create_lesson_config(self, lesson_data):
        """Tạo file config cho bài học"""
        config_file = self.output_dir / f"{lesson_data['lesson_id']}_config.json"
        
        with open(config_file, 'w', encoding='utf-8') as f:
            json.dump(lesson_data, f, ensure_ascii=False, indent=2)
        
        print(f"✅ Đã tạo config: {config_file}")
        return str(config_file)


def main():
    """Main function"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Chuyển đổi tài liệu')
    parser.add_argument('input_file', help='File cần chuyển đổi (PDF/DOCX/PPTX)')
    parser.add_argument('--output', '-o', help='Thư mục output')
    parser.add_argument('--type', '-t', choices=['pdf', 'docx', 'pptx'], help='Loại file')
    
    args = parser.parse_args()
    
    converter = DocumentConverter(output_dir=args.output or 'converted')
    
    input_path = Path(args.input_file)
    
    if not input_path.exists():
        print(f"❌ File không tồn tại: {input_path}")
        return
    
    # Auto-detect type from extension
    file_type = args.type or input_path.suffix[1:].lower()
    
    if file_type == 'pdf':
        converter.convert_pdf_to_images(input_path)
    elif file_type == 'docx':
        converter.convert_docx_to_html(input_path)
    elif file_type == 'pptx':
        result = converter.convert_pptx_to_images(input_path)
        if result:
            print(f"\n📋 Kết quả:")
            print(f"  - Tổng số slide: {result['total_slides']}")
            print(f"  - File ghi chú: {result['notes_file']}")
    else:
        print(f"❌ Định dạng không được hỗ trợ: {file_type}")


if __name__ == '__main__':
    if len(sys.argv) > 1:
        main()
    else:
        print("Script chuyển đổi tài liệu")
        print("\nCài đặt thư viện cần thiết:")
        print("  pip install pdf2image python-docx python-pptx Pillow")
        print("\nSử dụng:")
        print("  python convert_documents.py <file> [options]")
        print("\nVí dụ:")
        print("  python convert_documents.py document.pdf")
        print("  python convert_documents.py document.docx --output output_folder")
        print("  python convert_documents.py presentation.pptx")


